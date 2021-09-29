import os,glob
import wand
from wand.image import Image
from pptx import Presentation
from pptx.util import Inches 

""" summary_line

 Keyword arguments: create_ppt,slides,add_images
 argument -- description-array_of_images,titles-subtitles,ppt_name,path
 Return: This function takes an array of ppt_images,titles-subtitles array and output ppt name as an input arguments and creates a PowerPoint presentation document
 along with title and Subtitle for each image.

 """
 
def create_ppt(ppt_images,Titles_Subtitles,ppt_name,path):

    if not ppt_images:
        print('Oops!!, There are no images to create ppt document. Please check for water marked ".jpg" extension  files')
        return 
    # create an instance of Presentation
    prs = Presentation()

    # variable used to index the Titles_Subtitles array
    index=0

    # loop through all the images present in the ppt_images array
    for img in ppt_images:

        # slide_layouts[1] is used for title and content
        slide_layout = prs.slide_layouts[1]

        #add_slide function creates the new slide with the given layout
        slide = prs.slides.add_slide(slide_layout)

        # create a title
        title= slide.shapes.title

        #placholder for subtitle
        subtitle = slide.shapes.placeholders[1]

        # Assigning a title
        title.text = Titles_Subtitles[index][0]

        #Creating a text frame
        tf = subtitle.text_frame

        # Assigning a subtitle
        tf.text = Titles_Subtitles[index][1]

        # top,left are used to position the image in the slide
        top= Inches(2.5)
        left= Inches(1)

        # add_picture function adds the image to the slide
        pic = slide.shapes.add_picture(img, left, top)

        index+=1
    
    # Saving the ppt
    prs.save(ppt_name)

    # Print statement used to indicate that the requested ppt is created
    print("Your {ppt_name} document is ready to use".format(ppt_name=ppt_name))

    #Remove all the "watermarked_image*.jpg" after creating the document
    for file in ppt_images:
        os.remove(path+file)

    

"""summary_line

Keyword arguments: watermark
argument -- description: logo file, array of images, location or path the images are stored
Return: WaterMark function takes in 1 logo file,array of images and path as input and watermarks all the images present in 
an image array with the logo provided in logo_file at the specified path 

"""

def waterMark(logo_file,image_file_array,path):
    
    if not image_file_array:
        print('Oops!!, There are no images to WaterMark. Please check for ".jpg" extension files')
        return 
    try:
        with Image(filename=logo_file) as logo:
            
            # transform function is used to resize the image while maintaining the aspect ratio of the logo image
            logo.transform(resize="84x30")

            # Variable used for image count
            image_number=1

            # loop through all the image files present in the image_file_array
            for file in image_file_array:
                with Image(filename=path+file) as watermark_image:
                    # tranform function used to resize the image while maintaining the aspect ratio of the image
                    watermark_image.transform(resize="300x300")

                    # watermarks the image with the given logo 
                    watermark_image.watermark(logo,0.33,left=2,top=2)

                    # Saving the water marked image
                    watermark_image.save(filename="watermarked_image"+str(image_number)+".jpg")
                    image_number+=1

            # Print statement used to indicate the end of watermarking
            print("Congratulations..!!, All images are watermarked with the given logo..!!")

    except wand.exceptions.BlobError as error:
        print(error)

   
    
"""

Assumptions: 
           1. The single logo file is used to watermark images hence the logo file is hardcoded
           2. The provided logo file should be in ".png" extension format
           3. The naming convention of images that needs to be water marked should be in "image*.jpg" extension format
           4. In order to create ppt the naming convention for water marked image files should be "watermarked_image*.jpg" extension
           5. The Titles_Subtitles array for water marked images should predefined

"""

# Driver code
if __name__ == '__main__':

    #Holds the logo image file
    logo_file="nike_black.png"

    #Gets the current working directory path using os module
    path=os.getcwd()+'/'

    # glob function reads all the files with "image*.jpg" extension and creates a list/array of images files

    image_file_array=glob.glob('image*.jpg')

    # Invokes the WaterMark function 
    waterMark(logo_file,image_file_array,path)

    # Holds all the watermarked images 
    ppt_images=glob.glob("watermarked_image*.jpg")

    #Holds the titles[0] and subtitles[1] of each page of the ppt
    Titles_Subtitles=[["Slide1","Pots"],["Slide2","Oranges"],["Slide3","Desktop"],["Slide4","Cafe"],["Slide5","Camera"]]

    #Holds the name of the ppt to be created 
    ppt_name ="indycium_images.pptx"

    # Invokes the create_ppt function
    create_ppt(ppt_images,Titles_Subtitles,ppt_name,path)

    # opens the ppt file automatically 
    os.startfile(ppt_name)